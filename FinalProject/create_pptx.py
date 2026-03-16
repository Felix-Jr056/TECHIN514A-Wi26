from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN = RGBColor(0x00, 0xE6, 0x76)
RED = RGBColor(0xFF, 0x4D, 0x4D)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)
DARK_CARD = RGBColor(0x25, 0x25, 0x42)


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()


def add_title_bar(slide, title_text):
    # Title text
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    # Accent underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_bullet(tf, text, level=0, size=18, color=LIGHT_GRAY, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = Pt(6)
    return p


def add_card(slide, left, top, width, height, title, bullets):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_CARD
    card.line.fill.background()
    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.space_after = Pt(8)
    for b in bullets:
        add_bullet(tf, b, size=14)


# ── Slide 1: Title ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

tb = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "ErgoCompass"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Smart Posture Tracking System"
p2.font.size = Pt(28)
p2.font.color.rgb = ACCENT
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(8)

tb2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(1))
tf2 = tb2.text_frame
p3 = tf2.paragraphs[0]
p3.text = "TECHIN514A — Winter 2026"
p3.font.size = Pt(20)
p3.font.color.rgb = LIGHT_GRAY
p3.alignment = PP_ALIGN.CENTER

p4 = tf2.add_paragraph()
p4.text = "Felix Yu"
p4.font.size = Pt(22)
p4.font.color.rgb = WHITE
p4.font.bold = True
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(12)

tb3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.6))
tf3 = tb3.text_frame
p5 = tf3.paragraphs[0]
p5.text = "GitHub: https://github.com/Felix-Jr056/TECHIN514A-Wi26"
p5.font.size = Pt(16)
p5.font.color.rgb = ACCENT
p5.alignment = PP_ALIGN.CENTER

# ── Slide 2: Problem ────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "The Problem")

tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Poor sitting posture is a silent epidemic"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE
p.space_after = Pt(20)

bullets = [
    "80% of adults experience back pain at some point — poor posture is a leading cause",
    "Office workers sit 8+ hours/day, often unaware of their posture degradation",
    "Existing solutions are expensive, bulky, or require constant phone interaction",
    "No real-time, ambient feedback system exists for desk workers",
    "Posture problems develop gradually — by the time you feel pain, damage is done",
]
for b in bullets:
    add_bullet(tf, "•  " + b, size=20, color=WHITE)

# ── Slide 3: Proposed Solution ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Proposed Solution")

tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(1))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "ErgoCompass — A two-device wireless posture tracking system"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = ACCENT
p.space_after = Pt(16)

add_card(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(4),
         "Sensor Device (Chair Pad)", [
             "Fabric pad attached to chair back",
             "3× FSR 406 pressure sensors (top / mid / bottom)",
             "ADS1115 16-bit ADC for precise readings",
             "XIAO ESP32-C3 microcontroller",
             "BLE server — transmits at 5 Hz",
             "500 mAh LiPo battery",
         ])

add_card(slide, Inches(7), Inches(2.5), Inches(5.5), Inches(4),
         "Display Device (Desk Unit)", [
             "X27 stepper motor gauge needle (posture score)",
             "3× NeoPixel LEDs (per-sensor feedback)",
             "Green = good posture, Red = bad posture",
             "One-button calibration for personalization",
             "XIAO ESP32-C3 as BLE client",
             "Ambient feedback — no phone needed",
         ])

# ── Slide 4: System Architecture ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "System Architecture")

# Sensor box
s_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(4.5), Inches(4.8))
s_box.fill.solid()
s_box.fill.fore_color.rgb = DARK_CARD
s_box.line.color.rgb = ACCENT
s_box.line.width = Pt(2)

tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(4), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "SENSOR DEVICE"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT

tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(4), Inches(4))
tf = tb.text_frame
tf.word_wrap = True
items = [
    "3× FSR 406 Sensors",
    "    ↓  (Analog voltage)",
    "ADS1115 ADC (I2C, GAIN_TWO)",
    "    ↓  (16-bit digital)",
    "XIAO ESP32-C3",
    "    • EMA Filter (α=0.1)",
    "    • Baseline subtraction",
    "    • BLE Server (5 Hz notify)",
]
for i, item in enumerate(items):
    pp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    pp.text = item
    pp.font.size = Pt(14)
    pp.font.color.rgb = WHITE
    pp.space_after = Pt(4)

# Arrow
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.5), Inches(3.8), Inches(2), Inches(0.8))
arrow.fill.solid()
arrow.fill.fore_color.rgb = ACCENT
arrow.line.fill.background()
tb = slide.shapes.add_textbox(Inches(5.5), Inches(3.5), Inches(2), Inches(0.4))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "BLE (6 bytes)"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Display box
d_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.8), Inches(4.8))
d_box.fill.solid()
d_box.fill.fore_color.rgb = DARK_CARD
d_box.line.color.rgb = GREEN
d_box.line.width = Pt(2)

tb = slide.shapes.add_textbox(Inches(8.0), Inches(1.9), Inches(4.5), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "DISPLAY DEVICE"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = GREEN

tb = slide.shapes.add_textbox(Inches(8.0), Inches(2.5), Inches(4.5), Inches(4))
tf = tb.text_frame
tf.word_wrap = True
items = [
    "XIAO ESP32-C3",
    "    • BLE Client (scan + connect)",
    "    • Posture scoring algorithm",
    "        ↓",
    "X27 Stepper Motor Gauge",
    "    • 315° sweep, posture score",
    "3× NeoPixel LEDs",
    "    • Per-sensor color feedback",
    "Calibration Button (D6)",
]
for i, item in enumerate(items):
    pp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    pp.text = item
    pp.font.size = Pt(14)
    pp.font.color.rgb = WHITE
    pp.space_after = Pt(4)

# ── Slide 5: Critical Components / Pictures ─────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Critical Components")

placeholders = [
    (Inches(0.8), Inches(1.5), Inches(3.7), Inches(2.5), "Sensor PCB"),
    (Inches(4.8), Inches(1.5), Inches(3.7), Inches(2.5), "Display PCB"),
    (Inches(8.8), Inches(1.5), Inches(3.7), Inches(2.5), "Assembled Sensor Pad"),
    (Inches(0.8), Inches(4.3), Inches(3.7), Inches(2.5), "Assembled Display Unit"),
    (Inches(4.8), Inches(4.3), Inches(3.7), Inches(2.5), "FSR 406 Sensors"),
    (Inches(8.8), Inches(4.3), Inches(3.7), Inches(2.5), "X27 Motor + NeoPixels"),
]
for left, top, w, h, label in placeholders:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_CARD
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1)
    tb = slide.shapes.add_textbox(left, top + h / 2 - Inches(0.3), w, Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"[ {label} — Insert Photo ]"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

# ── Slide 6: Signal Processing ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Signal Processing & Algorithm")

add_card(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.5),
         "1. Sensor-Side: ADC + Filtering", [
             "ADS1115 at GAIN_TWO (±2.048V, 16-bit resolution)",
             "Baseline subtraction at boot (10-sample average, no pressure)",
             "Noise floor: readings < 200 clamped to 0",
             "EMA filter: α=0.1 (heavy smoothing for slow posture changes)",
             "Output: 3× int16 values at 5 Hz over BLE",
         ])

add_card(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(2.5),
         "2. Display-Side: Posture Scoring", [
             "Ideal values: Top=28000, Mid=15000, Bot=0",
             "Per-channel deviation normalized to 0–1",
             "Weighted combination: 35% top + 40% mid + 25% bot",
             "Score 0.0 (good) → needle 300° | Score 1.0 (bad) → needle 240°",
             "LED: green at ideal, linear gradient to red at deviation",
         ])

add_card(slide, Inches(0.8), Inches(4.3), Inches(11.8), Inches(2.5),
         "3. Accuracy & Performance", [
             "FSR response: ~0–30000 ADC range under typical sitting pressure",
             "Baseline calibration eliminates per-sensor DC offset (~5000 raw counts on some channels)",
             "EMA smoothing removes >90% of high-frequency noise while preserving posture transitions",
             "BLE latency: <200ms end-to-end (sensor read → display update)",
             "LED + motor update rate: real-time on each BLE notification (5 Hz)",
         ])

# ── Slide 7: Battery Considerations ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Battery Considerations")

add_card(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(3),
         "Battery Specification", [
             "Type: 1S LiPo (Li-Polymer)",
             "Model: PKCell 503035",
             "Capacity: 500 mAh @ 3.7V nominal",
             "Integrated PCM (protection circuit module)",
             "Used on both Sensor and Display devices",
         ])

add_card(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(3),
         "Power Budget", [
             "ESP32-C3: ~80mA active (BLE TX/RX)",
             "ADS1115: ~0.2mA continuous conversion",
             "FSR sensors: passive (no power draw)",
             "X27 motor: ~20mA when stepping",
             "NeoPixels: ~5mA × 3 at brightness 50/255",
         ])

add_card(slide, Inches(0.8), Inches(4.8), Inches(11.8), Inches(2),
         "Estimated Battery Life", [
             "Sensor: ~160mA typical → ~3 hours continuous use",
             "Display: ~120mA typical → ~4 hours continuous use",
             "Detailed analysis in Power Model spreadsheet (see Battery Life Estimate.pdf)",
         ])

# placeholder for chart
tb = slide.shapes.add_textbox(Inches(4), Inches(6.2), Inches(5), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "[ Insert battery life chart from Battery Life Estimate.pdf ]"
p.font.size = Pt(14)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.CENTER

# ── Slide 8: Budget Summary ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Budget Summary")

# Table
rows, cols = 12, 4
tbl_shape = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(1.5), Inches(10), Inches(5))
tbl = tbl_shape.table

tbl.columns[0].width = Inches(4)
tbl.columns[1].width = Inches(2.5)
tbl.columns[2].width = Inches(1.5)
tbl.columns[3].width = Inches(2)

data = [
    ["Component", "Model", "Qty", "Est. Cost"],
    ["Microcontroller", "XIAO ESP32-C3", "2", "$10.00"],
    ["ADC", "TI ADS1115IDGS", "1", "$3.00"],
    ["Pressure Sensors", "Interlink FSR 406", "3", "$18.00"],
    ["Stepper Motor", "X27", "1", "$5.00"],
    ["Motor Driver", "ULN2003A", "1", "$1.00"],
    ["NeoPixel LEDs", "SKC6812RV", "3", "$3.00"],
    ["Calibration Button", "Omron B3U-1000P", "1", "$0.50"],
    ["LiPo Batteries", "503035 500mAh", "2", "$8.00"],
    ["Custom PCBs", "JLCPCB fabrication", "2", "$10.00"],
    ["Misc (resistors, connectors)", "0805 / headers", "—", "$3.00"],
    ["", "", "TOTAL", "$61.50"],
]

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            if r == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = DARK_BG
            elif r == rows - 1:
                paragraph.font.bold = True
                paragraph.font.color.rgb = DARK_BG
            else:
                paragraph.font.color.rgb = DARK_BG
            if c >= 2:
                paragraph.alignment = PP_ALIGN.CENTER

# ── Slide 9: Future Work ────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Future Work")

tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
tf = tb.text_frame
tf.word_wrap = True

items = [
    ("Deep Sleep & Power Optimization", "Implement ESP32-C3 light-sleep between BLE intervals to extend battery life to 8+ hours"),
    ("Mobile App Integration", "Companion app for posture history, trends, and configurable alerts over BLE"),
    ("Machine Learning", "On-device posture classification (good/slouch/lean) trained on collected FSR data"),
    ("Haptic Feedback", "Add vibration motor to sensor pad for immediate tactile posture correction alerts"),
    ("Enclosure Design", "3D-printed cases for both devices — desk-friendly display and slim chair pad"),
    ("Multi-User Calibration", "Store multiple user profiles in flash for shared office chairs"),
]

for title, desc in items:
    p = tf.add_paragraph()
    p.text = f"▸  {title}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.space_before = Pt(12)

    p2 = tf.add_paragraph()
    p2.text = f"    {desc}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = LIGHT_GRAY
    p2.space_after = Pt(4)

# ── Slide 10: Thank You ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Questions?"
p2.font.size = Pt(30)
p2.font.color.rgb = ACCENT
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(16)

tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
tf2 = tb2.text_frame
p3 = tf2.paragraphs[0]
p3.text = "GitHub: https://github.com/Felix-Jr056/TECHIN514A-Wi26"
p3.font.size = Pt(18)
p3.font.color.rgb = LIGHT_GRAY
p3.alignment = PP_ALIGN.CENTER

# Save
output_path = "/Users/dyx/Documents/TECHIN514A-Wi26/FinalProject/ErgoCompass_Presentation.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
